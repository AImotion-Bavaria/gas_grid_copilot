from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load your existing ppt
ppt = Presentation('test.pptx')

# Select slides (0-indexed)
slide = ppt.slides[0]

# Modify shape properties -- for this example, we'll modify the first shape
shape = slide.shapes[0]

# Change fill color
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# Change line color
line = shape.line
line.color.rgb = RGBColor(0, 255, 0)  # Green color

# Change line width
line.width = Pt(2.0)

# Change text within the shape
if shape.has_text_frame:
    shape.text_frame.text = 'New text here'
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)

# Save ppt
ppt.save('Modified_Template.pptx')